# Create a presentation with one slide
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.util import Inches, Pt
import Excel_report

# whatis = lambda obj: print(type(obj), "\n\t" + "\n\t".join(dir(obj)))

# create presentation with 1 slide ------

def prepare_the_data(ws, column):
    list_data = []
    for row in range(2, ws.max_row):  # цикл по строкам с данными
        list_data.append(ws.cell(row, column).value)
    return list_data


def create_chart_line(ws, slide_num, column_1, column_2, column_3):
    chart_data = ChartData()
    chart_data.categories = prepare_the_data(ws, 1)
    if column_3 == 0 and column_2 != 0:
        chart_data.add_series(ws.cell(1, column_1).value, prepare_the_data(ws, column_1))
        chart_data.add_series(ws.cell(1, column_2).value, prepare_the_data(ws, column_2))
    elif column_3 == 0 and column_2 == 0:
        chart_data.add_series(ws.cell(1, column_1).value, prepare_the_data(ws, column_1))
    else:
        chart_data.add_series(ws.cell(1, column_1).value, prepare_the_data(ws, column_1))
        chart_data.add_series(ws.cell(1, column_2).value, prepare_the_data(ws, column_2))
        chart_data.add_series(ws.cell(1, column_3).value, prepare_the_data(ws, column_3))

    x, y, cx, cy = Inches(0.1), Inches(1), Inches(9.85), Inches(4.5)
    chart = slide_num.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.series[0].smooth = True


def create_chart_pie(ws, slide_num):
    chart_data = ChartData()
    chart_data.categories = [ws.cell(1, 2).value, ws.cell(1, 3).value]
    chart_data.add_series("False", (ws.cell(ws.max_row, 2).value, ws.cell(ws.max_row, 3).value))

    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(4.5)
    chart = slide_num.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0.0'
    data_labels.position = XL_LABEL_POSITION.OUTSIDE_END


def create_chart_column(ws, slide_num):
    chart_data = ChartData()
    chart_data.categories = prepare_the_data(ws, 1)
    chart_data.add_series(ws.cell(1, 2).value, prepare_the_data(ws, 2))
    chart_data.add_series(ws.cell(1, 3).value, prepare_the_data(ws, 3))
    chart_data.add_series(ws.cell(1, 4).value, prepare_the_data(ws, 4))
    chart_data.add_series(ws.cell(1, 5).value, prepare_the_data(ws, 5))
    chart_data.add_series(ws.cell(1, 6).value, prepare_the_data(ws, 6))
    chart_data.add_series(ws.cell(1, 7).value, prepare_the_data(ws, 7))
    chart_data.add_series(ws.cell(1, 8).value, prepare_the_data(ws, 8))

    x, y, cx, cy = Inches(0.1), Inches(1), Inches(9.85), Inches(4.5)
    graphic_frame = slide_num.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)

    chart.plots[0].has_data_labels = False


def sum_column_value(ws, column_num):
    num = 0
    for row in range(2, ws.max_row):  # цикл по строкам с данными
        num += ws.cell(row, column_num).value
    return str(round(num, 1))

def create_table(ws, slide_table):
    # задание параметров будущей таблицы
    rows = 11
    cols = 3
    left = Inches(0.2)
    top = Inches(1)
    width = Inches(10.0)
    height = Inches(4)

    table = slide_table.shapes.add_table(rows, cols, left, top, width, height).table  # создание таблицы
    # задание ширины столбцов
    table.columns[0].width = Inches(5)
    table.columns[1].width = Inches(1.9)
    table.columns[2].width = Inches(2.8)
    # write head cells
    table.cell(0, 0).text = 'Наименование параметра, единица измерения'
    table.cell(0, 1).text = 'Всего с начала\n2-ого этапа ОПЭ'
    table.cell(0, 2).text = 'Примечание'

    # write body cells
    table.cell(1, 0).text = 'Выработка электроэнергии БКЭУш, всего - кВт*ч'
    table.cell(2, 0).text = 'Выработка электроэнергии СМ – кВт*ч'
    table.cell(3, 0).text = 'Выработка электроэнергии ГПЭГ – кВт*ч'
    table.cell(4, 0).text = 'Расход электроэнергии на СН БКЭУш, кВт*ч'
    table.cell(5, 0).text = 'Расход электроэнергии на нужды площадки РН-21-22, кВт*ч'
    table.cell(6, 0).text = 'Выработка тепловой энергии от ГВК – кВт*ч'
    table.cell(7, 0).text = 'Наработка ГПЭГ - моточасы'
    table.cell(8, 0).text = 'Расход газа на ГВК – н.куб.м.'
    table.cell(9, 0).text = 'Расход газа на ГПЭГ – н.куб.м'
    table.cell(10, 0).text = 'Расход газа на БКЭУш, всего – н.куб.м'
    table.cell(7, 2).text = 'Последнее ТО – 26.01.21 – 230 мч'

    table.cell(1, 1).text = sum_column_value(ws, 4)
    table.cell(2, 1).text = sum_column_value(ws, 3)
    table.cell(3, 1).text = sum_column_value(ws, 2)
    table.cell(4, 1).text = sum_column_value(ws, 6)
    table.cell(5, 1).text = sum_column_value(ws, 5)
    table.cell(6, 1).text = sum_column_value(ws, 9)
    table.cell(7, 1).text = sum_column_value(ws, 7)
    table.cell(8, 1).text = sum_column_value(ws, 11)
    table.cell(9, 1).text = sum_column_value(ws, 10)
    table.cell(10, 1).text = sum_column_value(ws, 8)

    # Формат шрифта ячеек
    for cell in table.iter_cells():
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.name = "Franklin Gothic Book"
    # Формат шрифта заголовков ячеек первой строки
    table.cell(0, 0).text_frame.paragraphs[0].font.size = Pt(14)
    table.cell(0, 1).text_frame.paragraphs[1].font.size = Pt(14)
    table.cell(0, 1).text_frame.paragraphs[0].font.size = Pt(14)
    table.cell(0, 2).text_frame.paragraphs[0].font.size = Pt(14)
    table.cell(0, 0).text_frame.paragraphs[0].font.bold = True
    table.cell(0, 1).text_frame.paragraphs[1].font.bold = True
    table.cell(0, 1).text_frame.paragraphs[0].font.bold = True
    table.cell(0, 2).text_frame.paragraphs[0].font.bold = True
    table.cell(0, 1).text_frame.paragraphs[1].font.name = "Franklin Gothic Book"
    # Смещение данных относительно левого края ячейки
    table.cell(0, 0).margin_left = Inches(0.4)
    table.cell(0, 1).margin_left = Inches(0.3)
    table.cell(0, 2).margin_left = Inches(0.7)
    table.cell(0, 0).margin_top = Inches(0.15)
    table.cell(0, 2).margin_top = Inches(0.15)
    table.cell(1, 1).margin_left = Inches(0.7)
    table.cell(2, 1).margin_left = Inches(0.7)
    table.cell(3, 1).margin_left = Inches(0.7)
    table.cell(4, 1).margin_left = Inches(0.7)
    table.cell(5, 1).margin_left = Inches(0.7)
    table.cell(6, 1).margin_left = Inches(0.7)
    table.cell(7, 1).margin_left = Inches(0.7)
    table.cell(8, 1).margin_left = Inches(0.7)
    table.cell(9, 1).margin_left = Inches(0.7)
    table.cell(10, 1).margin_left = Inches(0.7)


def make_presentations(ws):

    prs = Presentation("temp_prs.pptx")
    blank_slide_chart = prs.slide_layouts[0]
    blank_slide_table = prs.slide_layouts[0]
    blank_slide_name = prs.slide_layouts[1]

    title_name = "Итоги ОПЭ БКЭУ-СМ/ГПЭГ/ГВК в ООО «Газпром добыча Кузнецк»\nна площадке РН-21-22\nза период с {} по {} г.".format(
        ws.cell(2, 15).value.strftime('%d.%m.%Y'),
        ws.cell(ws.max_row - 1, 16).value.strftime('%d.%m.%Y'))  # создание наименования заголовка заглавного слайда с автоматической вставкой дат начала и конца текущего ОПЭ
    title_table = "СВОДНЫЙ ОТЧЕТ"  # создание наименования заголовка слайда с таблицей
    # создание наименований заголовков слайдов с картинками
    title_0 = "ДИАГРАММА ИЗМЕРЯЕМЫХ ПРАМЕТРОВ ПО НЕДЕЛЯМ"
    title_1 = "ВЫРАБОТКА ЭЛЕКТРОЭНЕРГИИ"
    title_2 = "ПОТРЕБЛЕНИЕ ЭЛЕКТРОЭНЕРГИИ"
    title_3 = "ПОТРЕБЛЕНИЕ ГАЗА"
    title_4 = "ЭФФЕКТИВНОСТЬ ВЫРАБОТКИ ЭЛЕКТРОЭНЕРГИИ БКЭУ"
    title_5 = "ВЫРАБОТКА ТЕПЛОВОЙ ЭНЕРГИИ"
    title_6 = "ВЫРАБОТКА ЗА ОТЧЕТНЫЙ ПЕРИОД"
    title_7 = "ТЕМПЕРАТУРА ВНУТРИ БКЭУ"

    title_all = [title_0, title_1, title_2, title_3, title_4, title_5, title_6,
                 title_7]  # создание списка заголовков слайдов с картинками

    slide_name = prs.slides.add_slide(blank_slide_name)
    slide_table = prs.slides.add_slide(blank_slide_table)
    slide_0 = prs.slides.add_slide(blank_slide_chart)  # создание слайда с картинкой
    slide_1 = prs.slides.add_slide(blank_slide_chart)  # создание слайда с картинкой
    slide_2 = prs.slides.add_slide(blank_slide_chart)  # создание слайда с картинкой
    slide_3 = prs.slides.add_slide(blank_slide_chart)  # создание слайда с картинкой
    slide_4 = prs.slides.add_slide(blank_slide_chart)  # создание слайда с картинкой
    slide_5 = prs.slides.add_slide(blank_slide_chart)  # создание слайда с картинкой
    slide_6 = prs.slides.add_slide(blank_slide_chart)  # создание слайда с картинкой
    slide_7 = prs.slides.add_slide(blank_slide_chart)  # создание слайда с картинкой

    slide_all = [slide_0, slide_1, slide_2, slide_3, slide_4, slide_5, slide_6,
                 slide_7]  # создание списка слайдов с картинками

    title = slide_name.shapes.title  # создание заголовка заглавного слайда презентации
    title.text = title_name
    title = slide_table.shapes.title  # создание заголовка слайда презентации с таблицей
    title.text = title_table
    for i in range(0, len(slide_all)):  # создание заголовков слайдов презентаций с картинками
        title = slide_all[i].shapes.title
        title.text = title_all[i]


    create_table(ws, slide_table)
    create_chart_column(ws, slide_0)
    create_chart_line(ws, slide_1, 2, 3, 4)  # создание графика слайда "ДИАГРАММА ИЗМЕРЯЕМЫХ ПРАМЕТРОВ ПО НЕДЕЛЯМ"
    create_chart_line(ws, slide_2, 4, 5, 6)  # создание графика слайда "ВЫРАБОТКА ЭЛЕКТРОЭНЕРГИИ"
    create_chart_line(ws, slide_3, 8, 10, 11)  # создание графика слайда "ПОТРЕБЛЕНИЕ ГАЗА"
    create_chart_line(ws, slide_4, 12, 0, 0)  # создание графика слайда "ЭФФЕКТИВНОСТЬ ВЫРАБОТКИ ЭЛЕКТРОЭНЕРГИИ БКЭУ"
    create_chart_line(ws, slide_5, 9, 0, 0)  # создание графика слайда "ВЫРАБОТКА ТЕПЛОВОЙ ЭНЕРГИИ"
    create_chart_pie(ws, slide_6)  # создание диаграммы слайда "ВЫРАБОТКА ЗА ОТЧЕТНЫЙ ПЕРИОД"
    create_chart_line(ws, slide_7, 13, 14, 0)  # создание графика слайда "ТЕМПЕРАТУРА ВНУТРИ БКЭУ"

    prs.save("C:" + Excel_report.os.path.join(Excel_report.os.environ['HOMEPATH'],
                                              'Desktop') + '\\Промежуточные итоги ОПЭ БКЭУ.pptx')  # сохранение презентации на рабочем столе